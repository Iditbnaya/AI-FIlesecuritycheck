import os
import openai
from flask import Flask, request, jsonify, send_from_directory,session
from flask_cors import CORS
from pypdf import PdfReader
from dotenv import load_dotenv
from flask import Flask, send_from_directory
import xml.etree.ElementTree as ET  # Add this import for XML parsing
from docx import Document
from bs4 import BeautifulSoup
import json
import openpyxl
from pptx import Presentation
import csv

load_dotenv()


app = Flask(__name__)
CORS(app) 

# Load environment variables
endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
model_name = os.getenv("AZURE_OPENAI_DEPLOYMENT")
subscription_key = os.getenv("AZURE_OPENAI_KEY")
api_version = "2023-12-01-preview"

client = openai.AzureOpenAI(
    api_version=api_version,
    azure_endpoint=endpoint,
    api_key=subscription_key,
)

uploaded_text = "" # Variable to store the uploaded text
@app.route('/')
def serve_index():
    """Serve the main index.html file."""
    return send_from_directory('../frontend', 'index.html')

@app.errorhandler(404)
def not_found_error(error):
    """Handle 404 errors with JSON response."""
    return jsonify({"error": "Endpoint not found"}), 404

@app.errorhandler(500)
def internal_error(error):
    """Handle 500 errors with JSON response."""
    return jsonify({"error": "Internal server error"}), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    """ PDF, text, Office, or XML file upload """
    global uploaded_text
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded. Ensure the 'file' field is included in the request."}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file. Please choose a file to upload."}), 400

    try:
        filename = file.filename.lower()
        ext = os.path.splitext(filename)[1]
        file.seek(0)
        if ext == ".pdf":
            reader = PdfReader(file)
            uploaded_text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext == ".xml":
            try:
                tree = ET.parse(file)
                root = tree.getroot()
                uploaded_text = ET.tostring(root, encoding='unicode')
            except Exception:
                file.seek(0)
                uploaded_text = file.read().decode('utf-8', errors='ignore')
        elif ext == ".docx":
            file.seek(0)
            doc = Document(file)
            uploaded_text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in [".xls", ".xlsx"]:
            file.seek(0)
            wb = openpyxl.load_workbook(file, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            uploaded_text = text
        elif ext == ".csv":
            file.seek(0)
            reader = csv.reader(file.read().decode("utf-8", errors="ignore").splitlines())
            uploaded_text = "\n".join([", ".join(row) for row in reader])
        elif ext == ".pptx":
            file.seek(0)
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            uploaded_text = text
        else:
            file.seek(0)
            uploaded_text = file.read().decode('utf-8', errors='ignore')
    except ET.ParseError as e:
        return jsonify({"error": f"Failed to parse XML file: {str(e)}"}), 400
    except Exception as e:
        return jsonify({"error": f"Failed to process the file: {str(e)}"}), 500

    return jsonify({"message": "File uploaded successfully!"})

# Ensure static files and scripts are served correctly
@app.route('/scripts/<path:path>')
def serve_scripts(path):
    """Serve JavaScript files."""
    return send_from_directory('../frontend/scripts', path)

@app.route('/styles/<path:path>')
def serve_styles(path):
    """Serve CSS files."""
    return send_from_directory('../frontend/styles', path)

@app.route('/styles.css')
def serve_main_css():
    """Serve the main styles.css file from the frontend folder."""
    return send_from_directory('../frontend', 'styles.css')

@app.route('/ask', methods=['POST'])
def ask_question():
    """ Ask a question based on the uploaded text """
    global uploaded_text
    data = request.json
    question = data.get("question")

    if not uploaded_text:
        return jsonify({"error": "No file uploaded yet! Please upload a file first."}), 400
    if not question:
        return jsonify({"error": "No question provided! Please include a question in the request."}), 400

    try:
        prompt = f"""Based on the following document, answer the question below:\n\n{uploaded_text[:4000]}\n\nQuestion: {question}"""
        response = client.chat.completions.create(
            model=os.getenv("AZURE_OPENAI_DEPLOYMENT"),
            messages=[
                {"role": "system", "content": "You are an AI assistant."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=500
        )
        return jsonify({"answer": response.choices[0].message.content})
    except Exception as e:
        return jsonify({"error": f"Failed to process the question: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(debug=True)
